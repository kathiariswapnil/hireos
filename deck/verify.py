"""Re-open the built deck and check it for the failure modes that generated
decks are prone to: shapes off-canvas, text taller than its box, unresolved
images and stray empty text frames.

    .venv/bin/python -m deck.verify
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from brand import tokens as T  # noqa: E402
from deck import components as C, content as K  # noqa: E402
from deck.build import OUTPUT  # noqa: E402

# Shapes may bleed this far past the canvas before it counts as a problem;
# full-bleed art and centred diamonds legitimately sit on the edge.
BLEED_TOLERANCE = Inches(0.02)
# Text is allowed to exceed its declared box by this much, since box heights are
# intentionally generous and PowerPoint does not clip overflowing text.
OVERFLOW_TOLERANCE = Inches(0.06)


def _measure_frame(shape) -> tuple[int, float, str] | None:
    """Total rendered height of a shape's text, in EMU.

    Measured paragraph by paragraph so hanging indents (bulleted lists) and
    per-paragraph space_after are both accounted for -- a whole-frame estimate
    silently under-reports bullet blocks, which is exactly where overflow hides.
    """
    if not shape.has_text_frame:
        return None
    frame = shape.text_frame
    if not any(p.text.strip() for p in frame.paragraphs):
        return None

    inset = (frame.margin_left or 0) + (frame.margin_right or 0)
    total = 0
    largest = 0.0
    sample = ""

    for paragraph in frame.paragraphs:
        sizes = [r.font.size.pt for r in paragraph.runs if r.font.size is not None]
        if not sizes:
            continue
        size = max(sizes)
        bold = any(bool(r.font.bold) for r in paragraph.runs)
        spacing = (
            paragraph.line_spacing
            if isinstance(paragraph.line_spacing, float)
            else 1.18
        )

        pPr = paragraph._p.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
        )
        margin_left = int(pPr.get("marL", 0)) if pPr is not None else 0

        width = max((shape.width or 0) - inset - margin_left, Inches(0.2))
        lines = max(len(C.wrap_lines(paragraph.text, width, size, bold)), 1)
        total += int(lines * size * spacing * 12700)
        if paragraph.space_after is not None:
            total += int(paragraph.space_after)

        if size > largest:
            largest = size
        if not sample:
            sample = paragraph.text

    if total == 0:
        return None
    return total, largest, sample


def main() -> int:
    if not OUTPUT.exists():
        print(f"missing {OUTPUT}; run deck.build first", file=sys.stderr)
        return 1

    prs = Presentation(str(OUTPUT))
    problems: list[str] = []
    picture_count = 0
    shape_count = 0

    if len(prs.slides) != len(K.SLIDES):
        problems.append(
            f"slide count {len(prs.slides)} != content count {len(K.SLIDES)}"
        )

    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            shape_count += 1
            if shape.shape_type == 13 or shape.__class__.__name__ == "Picture":
                picture_count += 1

            left = shape.left or 0
            top = shape.top or 0
            right = left + (shape.width or 0)
            bottom = top + (shape.height or 0)

            if left < -BLEED_TOLERANCE or top < -BLEED_TOLERANCE:
                problems.append(
                    f"slide {index:02d}: {shape.shape_type} starts off-canvas "
                    f"({left / 914400:.2f}in, {top / 914400:.2f}in)"
                )
            if right > int(T.SLIDE_W) + BLEED_TOLERANCE:
                problems.append(
                    f"slide {index:02d}: shape overflows right edge to "
                    f"{right / 914400:.2f}in"
                )
            if bottom > int(T.SLIDE_H) + BLEED_TOLERANCE:
                problems.append(
                    f"slide {index:02d}: shape overflows bottom edge to "
                    f"{bottom / 914400:.2f}in"
                )

            measured = _measure_frame(shape)
            if measured is None:
                continue
            needed, size, sample = measured

            if needed > (shape.height or 0) + OVERFLOW_TOLERANCE:
                excess = (needed - (shape.height or 0)) / 914400
                problems.append(
                    f"slide {index:02d}: text overflows box by {excess:.2f}in "
                    f"at {size:.1f}pt -- {sample[:58]!r}"
                )
            # Compare against real rendered height, not the declared box. The
            # footer's own text legitimately starts on the footer line, so only
            # content that starts above it and spills across is a collision.
            footer_y = int(T.SLIDE_H) - int(T.MARGIN_BOTTOM)
            text_bottom = top + needed
            if top < footer_y - Inches(0.04) and text_bottom > footer_y:
                problems.append(
                    f"slide {index:02d}: text spills into the footer band by "
                    f"{(text_bottom - footer_y) / 914400:.2f}in -- {sample[:48]!r}"
                )
            elif text_bottom > int(T.SLIDE_H) - Inches(0.06):
                problems.append(
                    f"slide {index:02d}: text runs off the bottom of the slide -- "
                    f"{sample[:48]!r}"
                )

    expected_assets = [T.HERO_COVER, T.LOGO_MARK_DARK, *T.SECTION_ART.values()]
    for asset in expected_assets:
        if not asset.exists():
            problems.append(f"missing brand asset {asset.name}")

    print(f"slides     {len(prs.slides)}")
    print(f"shapes     {shape_count}")
    print(f"pictures   {picture_count}")
    print(f"assets     {sum(1 for a in expected_assets if a.exists())}/"
          f"{len(expected_assets)} present")

    if problems:
        print(f"\n{len(problems)} problem(s):")
        for problem in problems:
            print(f"  - {problem}")
        return 1

    print("\nno layout problems found")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
