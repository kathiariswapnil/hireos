"""Assemble the HireOS pre-seed investor deck.

    .venv/bin/python -m deck.build

Copy lives in content.py, primitives in components.py, diagrams in
diagrams.py. This module only decides where things go on each slide.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from brand import tokens as T  # noqa: E402
from deck import components as C, content as K, diagrams as D  # noqa: E402

OUTPUT = T.OUT_DIR / "HireOS_PreSeed_Investor_Deck.pptx"

# Lowest Y that slide content may occupy before it collides with the footer.
CONTENT_BOTTOM = int(T.SLIDE_H) - int(T.MARGIN_BOTTOM) - Inches(0.30)


# ==========================================================================
# Cover and closing
# ==========================================================================

def render_cover(slide, data: dict) -> None:
    if T.HERO_COVER.exists():
        C.add_picture_cover(slide, T.HERO_COVER, 0, 0, T.SLIDE_W, T.SLIDE_H)
    C.scrim(slide, 0, 0, int(T.SLIDE_W * 0.72), T.SLIDE_H,
            angle=0.0, start_alpha=0.97, end_alpha=0.0)

    C.logo_lockup(slide, T.MARGIN_X, T.MARGIN_TOP, Inches(0.52), dark=True)

    y = Inches(2.62)
    C.add_text(slide, T.MARGIN_X, y, Inches(7.0), Inches(0.28),
               data["eyebrow"].upper(), size=T.SIZE_EYEBROW, color=T.CYAN, bold=True,
               letter_spacing=1.6)

    y += Inches(0.42)
    C.add_text(slide, T.MARGIN_X, y, Inches(8.0), Inches(1.10), data["title"],
               size=T.SIZE_DISPLAY, color=T.TEXT_HI, bold=True, line_spacing=1.0,
               letter_spacing=-1.6)

    y += Inches(1.06)
    C.add_text(slide, T.MARGIN_X, y, Inches(7.4), Inches(0.44), data["descriptor"],
               size=Pt(22), color=T.TEXT_MID, line_spacing=1.2)

    y += Inches(0.66)
    C.gradient_bar(slide, T.MARGIN_X, y, Inches(2.6), Inches(0.045), T.CYAN, T.INDIGO)
    y += Inches(0.28)
    C.add_text(slide, T.MARGIN_X, y, Inches(7.4), Inches(0.36), data["tagline"],
               size=Pt(17), color=T.CYAN, bold=True)

    C.add_text(slide, T.MARGIN_X, int(T.SLIDE_H) - int(T.MARGIN_BOTTOM) - Inches(0.28),
               Inches(7.0), Inches(0.24), data["footnote"],
               size=T.SIZE_SMALL, color=T.TEXT_LOW)


def render_closing(slide, data: dict) -> None:
    if T.HERO_COVER.exists():
        picture = C.add_picture_cover(slide, T.HERO_COVER, 0, 0, T.SLIDE_W, T.SLIDE_H)
        C.set_picture_transparency(picture, 0.45)
    C.scrim(slide, 0, 0, int(T.SLIDE_W * 0.80), T.SLIDE_H,
            angle=0.0, start_alpha=0.96, end_alpha=0.0)

    C.logo_lockup(slide, T.MARGIN_X, T.MARGIN_TOP, Inches(0.52), dark=True)

    y = Inches(2.06)
    C.add_text(slide, T.MARGIN_X, y, Inches(7.6), Inches(2.05), data["title"],
               size=Pt(40), color=T.TEXT_HI, bold=True, line_spacing=1.16,
               letter_spacing=-1.0)

    y += Inches(2.46)
    C.rect(slide, T.MARGIN_X, y, Inches(0.045), Inches(1.16), fill=T.CYAN)
    C.add_text(slide, int(T.MARGIN_X) + Inches(0.28), y, Inches(7.4), Inches(1.16),
               data["body"], size=Pt(13), color=T.TEXT_MID, line_spacing=1.34)

    C.add_text(slide, T.MARGIN_X, int(T.SLIDE_H) - int(T.MARGIN_BOTTOM) - Inches(0.34),
               Inches(8.0), Inches(0.28), data["cta"],
               size=Pt(12), color=T.CYAN, bold=True)


def render_divider(slide, data: dict) -> None:
    art = T.SECTION_ART.get(data["art"])
    if art and art.exists():
        C.add_picture_cover(slide, art, 0, 0, T.SLIDE_W, T.SLIDE_H)
    C.scrim(slide, 0, 0, int(T.SLIDE_W * 0.68), T.SLIDE_H,
            angle=0.0, start_alpha=0.95, end_alpha=0.0)

    C.add_text(slide, T.MARGIN_X, Inches(2.48), Inches(2.0), Inches(0.88),
               data["number"], size=Pt(56), color=C._mix(T.INK, T.INDIGO, 0.55),
               bold=True, line_spacing=1.0)
    C.gradient_bar(slide, T.MARGIN_X, Inches(3.54), Inches(1.7), Inches(0.045),
                   T.CYAN, T.INDIGO)
    C.add_text(slide, T.MARGIN_X, Inches(3.82), Inches(7.0), Inches(0.62),
               data["title"], size=Pt(40), color=T.TEXT_HI, bold=True,
               letter_spacing=-1.0)
    C.add_text(slide, T.MARGIN_X, Inches(4.54), Inches(6.4), Inches(0.60),
               data["subtitle"], size=Pt(15), color=T.TEXT_MID, line_spacing=1.3)


# ==========================================================================
# Generic card row
# ==========================================================================

def render_cards(slide, data: dict, y: int) -> None:
    cards = data["cards"]
    reserved = 0
    if data.get("kicker"):
        reserved += Inches(0.62)
    if data.get("disclaimer"):
        reserved += Inches(0.58)

    height = CONTENT_BOTTOM - y - reserved
    height = min(height, Inches(3.30))
    geo = T.columns(len(cards), gutter=Inches(0.26))

    for index, (spec, (cx, cw)) in enumerate(zip(cards, geo)):
        C.card(slide, cx, y, cw, height,
               heading=spec["heading"], body=spec["body"],
               accent=T.ACCENTS[index % len(T.ACCENTS)],
               number=f"{index + 1:02d}")

    y += int(height) + Inches(0.24)
    if data.get("kicker"):
        y += C.kicker(slide, T.MARGIN_X, y, Inches(9.0), data["kicker"]) + Inches(0.16)
    if data.get("disclaimer"):
        C.note_bar(slide, T.MARGIN_X, y, T.CONTENT_W, data["disclaimer"],
                   label="Note", accent=T.TEXT_LOW, size=T.SIZE_MICRO)


# ==========================================================================
# Candidate intelligence - evidence table
# ==========================================================================

_ASSESSMENT_COLORS = {
    "Very strong": T.GREEN,
    "Strong": T.CYAN,
    "Gap to validate": T.AMBER,
}


def render_evidence(slide, data: dict, y: int) -> None:
    table_w = Inches(7.55)
    right_x = int(T.MARGIN_X) + int(table_w) + Inches(0.44)
    right_w = int(T.MARGIN_X) + int(T.CONTENT_W) - right_x

    C.add_text(slide, T.MARGIN_X, y, table_w, Inches(0.26),
               data["candidate_label"].upper(), size=T.SIZE_MICRO, color=T.CYAN,
               bold=True, letter_spacing=1.2)
    ty = y + Inches(0.34)

    col_w = (Inches(1.72), Inches(4.05), Inches(1.78))
    header_h = Inches(0.34)
    C.rect(slide, T.MARGIN_X, ty, table_w, header_h, fill=T.SURFACE_2, line=T.LINE)
    cx = int(T.MARGIN_X)
    for label, width in zip(data["columns"], col_w):
        C.add_text(slide, cx + Inches(0.14), ty, int(width) - Inches(0.20), header_h,
                   label.upper(), size=Pt(8.5), color=T.TEXT_LOW, bold=True,
                   anchor=MSO_ANCHOR.MIDDLE, letter_spacing=1.0)
        cx += int(width)
    ty += int(header_h)

    row_h = Inches(0.56)
    for index, (competency, evidence, assessment) in enumerate(data["rows"]):
        accent = _ASSESSMENT_COLORS.get(assessment, T.TEXT_LOW)
        C.rect(slide, T.MARGIN_X, ty, table_w, row_h,
               fill=T.SURFACE if index % 2 == 0 else T.SURFACE_2, line=T.LINE_SOFT)
        cx = int(T.MARGIN_X)
        C.add_text(slide, cx + Inches(0.14), ty, int(col_w[0]) - Inches(0.20), row_h,
                   competency, size=Pt(11.5), color=T.TEXT_HI, bold=True,
                   anchor=MSO_ANCHOR.MIDDLE)
        cx += int(col_w[0])
        C.add_text(slide, cx + Inches(0.14), ty, int(col_w[1]) - Inches(0.24), row_h,
                   evidence, size=Pt(10.5), color=T.TEXT_MID, anchor=MSO_ANCHOR.MIDDLE,
                   line_spacing=1.18)
        cx += int(col_w[1])
        C.rect(slide, cx + Inches(0.14), ty + Inches(0.135),
               int(col_w[2]) - Inches(0.30), Inches(0.29),
               fill=C._mix(T.INK, accent, 0.18), line=C._mix(T.LINE, accent, 0.45),
               radius=0.22)
        C.add_text(slide, cx + Inches(0.14), ty + Inches(0.135),
                   int(col_w[2]) - Inches(0.30), Inches(0.29), assessment,
                   size=Pt(9.5), color=accent, bold=True, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE)
        ty += int(row_h)

    # Recommendation panel.
    ry = y + Inches(0.34)
    panel_h = Inches(1.66)
    C.rect(slide, right_x, ry, right_w, panel_h,
           fill=C._mix(T.SURFACE, T.GREEN, 0.10), line=C._mix(T.LINE, T.GREEN, 0.40),
           radius=0.06)
    C.rect(slide, right_x, ry, right_w, Inches(0.04), fill=T.GREEN)
    pad = Inches(0.24)
    C.add_text(slide, right_x + int(pad), ry + Inches(0.20), right_w - 2 * int(pad),
               Inches(0.24), data["recommendation_label"].upper(),
               size=T.SIZE_MICRO, color=T.GREEN, bold=True, letter_spacing=1.2)
    C.add_text(slide, right_x + int(pad), ry + Inches(0.46), right_w - 2 * int(pad),
               Inches(0.32), data["recommendation"],
               size=Pt(17), color=T.TEXT_HI, bold=True)
    C.add_text(slide, right_x + int(pad), ry + Inches(0.82), right_w - 2 * int(pad),
               Inches(0.82), data["reason"], size=Pt(10), color=T.TEXT_MID,
               line_spacing=1.28)

    oy = ry + int(panel_h) + Inches(0.20)
    C.add_text(slide, right_x, oy, right_w, Inches(0.24), "EVERY EVALUATION RETURNS",
               size=T.SIZE_MICRO, color=T.CYAN, bold=True, letter_spacing=1.2)
    oy += Inches(0.28)
    for label, description in data["outputs"]:
        C.diamond(slide, right_x + Inches(0.06), oy + Inches(0.09), Inches(0.13),
                  fill=T.CYAN)
        C.add_text(slide, right_x + Inches(0.22), oy, right_w - Inches(0.22),
                   Inches(0.20), label, size=Pt(10.5), color=T.TEXT_HI, bold=True)
        C.add_text(slide, right_x + Inches(0.22), oy + Inches(0.21),
                   right_w - Inches(0.22), Inches(0.24), description,
                   size=Pt(9.5), color=T.TEXT_MID, line_spacing=1.2)
        oy += Inches(0.44)


# ==========================================================================
# Autonomy ladder
# ==========================================================================

def render_autonomy(slide, data: dict, y: int) -> None:
    modes = data["modes"]
    geo = T.columns(len(modes), gutter=Inches(0.30))
    base_h = Inches(2.30)
    rise = Inches(0.30)
    accents = (T.TEXT_LOW, T.INDIGO, T.CYAN)
    bottom = y + int(base_h) + 2 * int(rise)

    for index, (spec, (cx, cw)) in enumerate(zip(modes, geo)):
        accent = accents[index]
        height = int(base_h) + index * int(rise)
        top = bottom - height
        is_core = index == 1

        C.rect(slide, cx, top, cw, height,
               fill=C._mix(T.SURFACE, accent, 0.14) if is_core else T.SURFACE,
               line=accent if is_core else T.LINE, radius=0.05)
        C.rect(slide, cx, top, cw, Inches(0.045), fill=accent)

        pad = Inches(0.28)
        inner_x = int(cx) + int(pad)
        inner_w = int(cw) - 2 * int(pad)
        ty = top + int(Inches(0.26))

        C.add_text(slide, inner_x, ty, inner_w, Inches(0.24),
                   f"RUNG {index + 1}", size=Pt(8.5), color=accent, bold=True,
                   letter_spacing=1.3)
        ty += Inches(0.28)
        C.add_text(slide, inner_x, ty, inner_w, Inches(0.40), spec["name"],
                   size=Pt(24), color=T.TEXT_HI, bold=True, letter_spacing=-0.6)
        ty += Inches(0.50)
        C.add_text(slide, inner_x, ty, inner_w, Inches(0.52), spec["behavior"],
                   size=Pt(12.5), color=T.TEXT_HI, bold=True, line_spacing=1.24)
        ty += Inches(0.58)
        C.add_text(slide, inner_x, ty, inner_w, Inches(0.24), "EXAMPLE",
                   size=Pt(8), color=T.TEXT_LOW, bold=True, letter_spacing=1.2)
        ty += Inches(0.22)
        C.add_text(slide, inner_x, ty, inner_w, Inches(0.46), spec["example"],
                   size=Pt(10.5), color=T.TEXT_MID, line_spacing=1.26)

        C.rect(slide, inner_x, bottom - Inches(0.52), inner_w, Inches(0.32),
               fill=C._mix(T.INK, accent, 0.20), line=C._mix(T.LINE, accent, 0.45),
               radius=0.20)
        C.add_text(slide, inner_x, bottom - Inches(0.52), inner_w, Inches(0.32),
                   spec["trust"], size=Pt(9.5), color=accent, bold=True,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Rising trust axis beneath the rungs.
    ay = bottom + Inches(0.20)
    C.arrow(slide, T.MARGIN_X, ay, int(T.MARGIN_X) + int(T.CONTENT_W), ay,
            T.LINE, width=Pt(1.25))
    C.add_text(slide, T.MARGIN_X, ay + Inches(0.08), Inches(6.0), Inches(0.24),
               "INCREASING DELEGATED AUTHORITY, EARNED WITH EVIDENCE",
               size=Pt(8.5), color=T.TEXT_LOW, bold=True, letter_spacing=1.2)

    C.kicker(slide, T.MARGIN_X, ay + Inches(0.42), Inches(10.5), data["kicker"])


# ==========================================================================
# Wedge and MVP
# ==========================================================================

def render_wedge(slide, data: dict, y: int) -> None:
    steps = data["wedge_steps"]
    overlap = Inches(0.11)
    stage_w = int((int(T.CONTENT_W) + (len(steps) - 1) * int(overlap)) / len(steps))
    stage_h = Inches(0.50)

    C.add_text(slide, T.MARGIN_X, y, Inches(4.0), Inches(0.24),
               data["wedge_label"].upper(), size=T.SIZE_MICRO, color=T.CYAN, bold=True,
               letter_spacing=1.2)
    y += Inches(0.30)

    for index, step in enumerate(steps):
        sx = int(T.MARGIN_X) + index * (stage_w - int(overlap))
        accent = C._mix(T.INDIGO, T.CYAN, index / max(len(steps) - 1, 1))
        shape = C.rect(slide, sx, y, stage_w, stage_h,
                       fill=C._mix(T.SURFACE, accent, 0.20), line=accent,
                       shape_type=MSO_SHAPE.CHEVRON)
        shape.line.width = Pt(0.75)
        size = C.fit_size(step, stage_w - Inches(0.38), stage_h, 11.0, bold=True,
                          line_spacing=1.05)
        C.add_text(slide, sx + Inches(0.19), y, stage_w - Inches(0.38), stage_h, step,
                   size=Pt(size), color=T.TEXT_HI, bold=True, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

    y += int(stage_h) + Inches(0.24)

    build_w = Inches(6.30)
    defer_x = int(T.MARGIN_X) + int(build_w) + Inches(0.42)
    defer_w = int(T.MARGIN_X) + int(T.CONTENT_W) - defer_x
    panel_h = Inches(1.96)

    for x, width, spec, accent, mark in (
        (int(T.MARGIN_X), int(build_w), data["build"], T.CYAN, "\u2022"),
        (defer_x, defer_w, data["defer"], T.ROSE, "\u2013"),
    ):
        C.rect(slide, x, y, width, panel_h, fill=T.SURFACE, line=T.LINE, radius=0.05)
        C.rect(slide, x, y, width, Inches(0.04), fill=accent)
        pad = Inches(0.26)
        C.add_text(slide, x + int(pad), y + Inches(0.20), width - 2 * int(pad),
                   Inches(0.26), spec["heading"], size=Pt(14), color=T.TEXT_HI, bold=True)

        items = spec["items"]
        half = (len(items) + 1) // 2
        sub_w = int((width - 2 * int(pad) - Inches(0.24)) / 2)
        for col_index, chunk in enumerate((items[:half], items[half:])):
            C.add_bullets(slide, x + int(pad) + col_index * (sub_w + int(Inches(0.24))),
                          y + Inches(0.52), sub_w, panel_h - Inches(0.68), chunk,
                          size=Pt(9.5), color=T.TEXT_MID, accent=accent,
                          line_spacing=1.16, space_after=Pt(3.5), bullet_char=mark)

    y += int(panel_h) + Inches(0.18)
    C.kicker(slide, T.MARGIN_X, y, Inches(10.8), data["kicker"])


# ==========================================================================
# Competitive matrix
# ==========================================================================

def render_matrix(slide, data: dict, y: int) -> None:
    height = Inches(3.62)
    C.styled_table(
        slide, T.MARGIN_X, y, T.CONTENT_W, height,
        list(data["columns"]), data["rows"],
        col_widths=[2.5, 1.55, 1.2, 1.2, 1.35],
        highlight_column=data.get("highlight_column"),
        header_size=Pt(10.5), body_size=Pt(10.5), row_height=Inches(0.40),
    )
    C.kicker(slide, T.MARGIN_X, y + int(height) + Inches(0.22), Inches(10.8),
             data["kicker"])


# ==========================================================================
# Moat
# ==========================================================================

def render_moat(slide, data: dict, y: int) -> None:
    layers = data["layers"]
    row_h = Inches(0.53)
    gap = Inches(0.07)

    for index, (heading, body) in enumerate(layers):
        ry = y + index * (int(row_h) + int(gap))
        accent = C._mix(T.INDIGO, T.CYAN, index / max(len(layers) - 1, 1))
        C.rect(slide, T.MARGIN_X, ry, T.CONTENT_W, row_h,
               fill=T.SURFACE if index % 2 == 0 else T.SURFACE_2,
               line=T.LINE_SOFT, radius=0.10)
        C.rect(slide, T.MARGIN_X, ry, Inches(0.04), row_h, fill=accent)
        C.add_text(slide, int(T.MARGIN_X) + Inches(0.22), ry, Inches(0.36), row_h,
                   f"{index + 1:02d}", size=Pt(9.5), color=accent, bold=True,
                   anchor=MSO_ANCHOR.MIDDLE)
        C.add_text(slide, int(T.MARGIN_X) + Inches(0.62), ry, Inches(3.55), row_h,
                   heading, size=Pt(12.5), color=T.TEXT_HI, bold=True,
                   anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)
        C.add_text(slide, int(T.MARGIN_X) + Inches(4.34), ry,
                   int(T.CONTENT_W) - Inches(4.56), row_h, body,
                   size=Pt(10.5), color=T.TEXT_MID, anchor=MSO_ANCHOR.MIDDLE,
                   line_spacing=1.2)

    y += len(layers) * (int(row_h) + int(gap)) + Inches(0.16)
    C.kicker(slide, T.MARGIN_X, y, Inches(10.8), data["kicker"])


# ==========================================================================
# Pricing
# ==========================================================================

def render_pricing(slide, data: dict, y: int) -> None:
    packages = data["packages"]
    geo = T.columns(len(packages), gutter=Inches(0.24))
    height = Inches(2.86)

    for index, (spec, (cx, cw)) in enumerate(zip(packages, geo)):
        highlight = spec.get("highlight")
        accent = T.CYAN if highlight else T.INDIGO
        C.rect(slide, cx, y, cw, height,
               fill=C._mix(T.SURFACE, accent, 0.12) if highlight else T.SURFACE,
               line=accent if highlight else T.LINE, radius=0.05)
        C.rect(slide, cx, y, cw, Inches(0.045), fill=accent)

        pad = Inches(0.26)
        inner_x = int(cx) + int(pad)
        inner_w = int(cw) - 2 * int(pad)
        ty = y + int(Inches(0.24))

        if spec.get("note"):
            width = C.chip(slide, inner_x, ty, Inches(0.26), spec["note"].upper(),
                           fill=accent, line=None, color=T.INK, size=Pt(8), bold=True)
            ty += Inches(0.34)
        else:
            ty += Inches(0.34)

        C.add_text(slide, inner_x, ty, inner_w, Inches(0.26), spec["name"],
                   size=Pt(13.5), color=T.TEXT_HI, bold=True)
        ty += Inches(0.34)
        size = C.fit_size(spec["price"], inner_w, Inches(0.40), 21.0, bold=True,
                          line_spacing=1.0)
        C.add_text(slide, inner_x, ty, inner_w, Inches(0.38), spec["price"],
                   size=Pt(size), color=accent, bold=True, letter_spacing=-0.5)
        ty += Inches(0.40)
        C.add_text(slide, inner_x, ty, inner_w, Inches(0.22),
                   spec["period"].upper() if spec["period"] else "",
                   size=Pt(8.5), color=T.TEXT_LOW, bold=True, letter_spacing=1.1)
        ty += Inches(0.30)
        C.rect(slide, inner_x, ty, inner_w, Emu(9525), fill=T.LINE)
        ty += Inches(0.14)
        C.add_bullets(slide, inner_x, ty, inner_w, y + int(height) - ty - Inches(0.16),
                      spec["includes"], size=Pt(10), accent=accent, space_after=Pt(4))

    y += int(height) + Inches(0.24)
    C.note_bar(slide, T.MARGIN_X, y, T.CONTENT_W, data["assumption"],
               label="Assumption", accent=T.AMBER)


# ==========================================================================
# Market
# ==========================================================================

def render_market(slide, data: dict, y: int) -> None:
    left_w = Inches(5.55)
    right_x = int(T.MARGIN_X) + int(left_w) + Inches(0.48)
    right_w = int(T.MARGIN_X) + int(T.CONTENT_W) - right_x

    C.add_text(slide, T.MARGIN_X, y, left_w, Inches(0.24),
               data["assumptions_label"].upper(), size=T.SIZE_MICRO, color=T.AMBER,
               bold=True, letter_spacing=1.2)
    ay = y + Inches(0.32)
    row_h = Inches(0.44)
    for index, (label, value) in enumerate(data["assumptions"]):
        C.rect(slide, T.MARGIN_X, ay, left_w, row_h,
               fill=T.SURFACE if index % 2 == 0 else T.SURFACE_2, line=T.LINE_SOFT)
        C.add_text(slide, int(T.MARGIN_X) + Inches(0.18), ay, Inches(2.10), row_h,
                   label, size=Pt(10), color=T.TEXT_LOW, anchor=MSO_ANCHOR.MIDDLE,
                   line_spacing=1.15)
        C.add_text(slide, int(T.MARGIN_X) + Inches(2.34), ay,
                   int(left_w) - Inches(2.52), row_h, value,
                   size=Pt(10.5), color=T.TEXT_HI, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                   line_spacing=1.15)
        ay += int(row_h)

    C.add_text(slide, right_x, y, right_w, Inches(0.24), "BOTTOM-UP ARITHMETIC",
               size=T.SIZE_MICRO, color=T.CYAN, bold=True, letter_spacing=1.2)
    ty = y + Inches(0.32)
    tier_h = Inches(0.72)
    accents = (T.INDIGO, T.CYAN, T.GREEN)
    for index, (label, value, math) in enumerate(data["tiers"]):
        accent = accents[index % len(accents)]
        C.rect(slide, right_x, ty, right_w, tier_h,
               fill=C._mix(T.SURFACE, accent, 0.10), line=C._mix(T.LINE, accent, 0.35),
               radius=0.06)
        C.rect(slide, right_x, ty, Inches(0.04), tier_h, fill=accent)
        C.add_text(slide, right_x + Inches(0.22), ty + Inches(0.10), Inches(2.4),
                   Inches(0.22), label.upper(), size=Pt(8.5), color=T.TEXT_LOW,
                   bold=True, letter_spacing=1.1)
        C.add_text(slide, right_x + Inches(0.22), ty + Inches(0.32), Inches(2.6),
                   Inches(0.34), value, size=Pt(20), color=accent, bold=True,
                   letter_spacing=-0.6)
        C.add_text(slide, right_x + Inches(2.95), ty, right_w - Inches(3.15), tier_h,
                   math, size=Pt(9.5), color=T.TEXT_MID, anchor=MSO_ANCHOR.MIDDLE,
                   align=PP_ALIGN.RIGHT, line_spacing=1.2)
        ty += int(tier_h) + Inches(0.12)

    y = max(ay, ty) + Inches(0.12)
    y += C.kicker(slide, T.MARGIN_X, y, Inches(10.8), data["reality_check"]) + Inches(0.16)
    C.note_bar(slide, T.MARGIN_X, y, T.CONTENT_W, data["caveat"],
               label="Caveat", accent=T.TEXT_LOW, size=T.SIZE_MICRO)


# ==========================================================================
# Validation plan
# ==========================================================================

def render_validation(slide, data: dict, y: int) -> None:
    left_w = Inches(3.70)
    right_x = int(T.MARGIN_X) + int(left_w) + Inches(0.48)
    right_w = int(T.MARGIN_X) + int(T.CONTENT_W) - right_x

    C.add_text(slide, T.MARGIN_X, y, left_w, Inches(0.24), "HONEST STATUS",
               size=T.SIZE_MICRO, color=T.ROSE, bold=True, letter_spacing=1.2)
    sy = y + Inches(0.32)
    for label, value in data["status"]:
        height = C.text_height(value, int(left_w) - Inches(0.36), 10.5,
                               line_spacing=1.22) + Inches(0.38)
        C.rect(slide, T.MARGIN_X, sy, left_w, height,
               fill=T.SURFACE, line=T.LINE, radius=0.07)
        C.rect(slide, T.MARGIN_X, sy, Inches(0.035), height, fill=T.ROSE)
        C.add_text(slide, int(T.MARGIN_X) + Inches(0.20), sy + Inches(0.11),
                   int(left_w) - Inches(0.36), Inches(0.20), label.upper(),
                   size=Pt(8.5), color=T.TEXT_LOW, bold=True, letter_spacing=1.1)
        C.add_text(slide, int(T.MARGIN_X) + Inches(0.20), sy + Inches(0.29),
                   int(left_w) - Inches(0.36), height - Inches(0.36), value,
                   size=Pt(10.5), color=T.TEXT_HI, line_spacing=1.22)
        sy += int(height) + Inches(0.08)

    C.add_text(slide, right_x, y, right_w, Inches(0.24),
               "HOW WE EARN CONVICTION", size=T.SIZE_MICRO, color=T.CYAN, bold=True,
               letter_spacing=1.2)
    ry = y + Inches(0.32)
    step_h = Inches(0.52)
    for number, heading, body in data["steps"]:
        C.rect(slide, right_x, ry, right_w, step_h, fill=T.SURFACE,
               line=T.LINE_SOFT, radius=0.07)
        C.rect(slide, right_x + Inches(0.14), ry + Inches(0.11), Inches(0.30),
               Inches(0.30), fill=T.INDIGO, radius=0.30)
        C.add_text(slide, right_x + Inches(0.14), ry + Inches(0.11), Inches(0.30),
                   Inches(0.30), number, size=Pt(10), color=T.TEXT_HI, bold=True,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        C.add_text(slide, right_x + Inches(0.56), ry + Inches(0.06), Inches(2.55),
                   Inches(0.40), heading, size=Pt(11), color=T.TEXT_HI, bold=True,
                   anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
        C.add_text(slide, right_x + Inches(3.20), ry, right_w - Inches(3.40), step_h,
                   body, size=Pt(9.5), color=T.TEXT_MID, anchor=MSO_ANCHOR.MIDDLE,
                   line_spacing=1.2)
        ry += int(step_h) + Inches(0.06)

    C.kicker(slide, T.MARGIN_X, max(sy, ry) + Inches(0.02), Inches(10.8),
             data["kicker"])


# ==========================================================================
# North-star metrics
# ==========================================================================

def render_metrics(slide, data: dict, y: int) -> None:
    star = data["north_star"]
    hero_h = Inches(1.52)
    C.rect(slide, T.MARGIN_X, y, T.CONTENT_W, hero_h,
           fill=C._mix(T.SURFACE, T.CYAN, 0.13), line=T.CYAN, radius=0.05)
    C.gradient_bar(slide, T.MARGIN_X, y, T.CONTENT_W, Inches(0.05), T.CYAN, T.INDIGO)

    C.add_text(slide, int(T.MARGIN_X) + Inches(0.32), y + Inches(0.24),
               Inches(3.0), Inches(0.24), star["label"].upper(),
               size=T.SIZE_MICRO, color=T.CYAN, bold=True, letter_spacing=1.3)
    C.add_text(slide, int(T.MARGIN_X) + Inches(0.32), y + Inches(0.52),
               Inches(7.6), Inches(0.46), star["metric"],
               size=Pt(28), color=T.TEXT_HI, bold=True, letter_spacing=-0.9)
    C.add_text(slide, int(T.MARGIN_X) + Inches(0.32), y + Inches(1.02),
               Inches(8.4), Inches(0.28), star["why"],
               size=Pt(11.5), color=T.TEXT_MID)
    C.diamond(slide, int(T.MARGIN_X) + int(T.CONTENT_W) - Inches(0.86),
              y + int(hero_h / 2), Inches(0.66), fill=None, line=T.CYAN, line_w=Pt(1.5))
    C.diamond(slide, int(T.MARGIN_X) + int(T.CONTENT_W) - Inches(0.86),
              y + int(hero_h / 2), Inches(0.28), fill=T.CYAN)

    y += int(hero_h) + Inches(0.28)
    C.add_text(slide, T.MARGIN_X, y, T.CONTENT_W, Inches(0.24), "SUPPORTING METRICS",
               size=T.SIZE_MICRO, color=T.TEXT_LOW, bold=True, letter_spacing=1.2)
    y += Inches(0.30)

    supporting = data["supporting"]
    cols = 4
    gap = Inches(0.20)
    card_w = int((int(T.CONTENT_W) - (cols - 1) * int(gap)) / cols)
    card_h = Inches(0.94)
    for index, (metric, why) in enumerate(supporting):
        row, col = divmod(index, cols)
        cx = int(T.MARGIN_X) + col * (card_w + int(gap))
        cy = y + row * (int(card_h) + int(Inches(0.12)))
        accent = T.ACCENTS[index % len(T.ACCENTS)]
        C.rect(slide, cx, cy, card_w, card_h, fill=T.SURFACE, line=T.LINE, radius=0.07)
        C.rect(slide, cx, cy, Inches(0.03), card_h, fill=accent)
        C.add_text(slide, cx + Inches(0.20), cy + Inches(0.13), card_w - Inches(0.36),
                   Inches(0.42), metric, size=Pt(11), color=T.TEXT_HI, bold=True,
                   line_spacing=1.14)
        C.add_text(slide, cx + Inches(0.20), cy + Inches(0.57), card_w - Inches(0.36),
                   Inches(0.30), why, size=Pt(9.5), color=T.TEXT_MID, line_spacing=1.2)


# ==========================================================================
# Roadmap
# ==========================================================================

def render_roadmap(slide, data: dict, y: int) -> None:
    sprints = data["sprints"]
    cols = 6
    gap = Inches(0.16)
    card_w = int((int(T.CONTENT_W) - (cols - 1) * int(gap)) / cols)

    C.add_text(slide, T.MARGIN_X, y, Inches(5.0), Inches(0.24), "90-DAY BUILD PLAN",
               size=T.SIZE_MICRO, color=T.CYAN, bold=True, letter_spacing=1.2)
    y += Inches(0.34)

    # Timeline spine.
    C.rect(slide, T.MARGIN_X, y + Inches(0.10), T.CONTENT_W, Emu(12700), fill=T.LINE)
    card_h = Inches(2.04)
    for index, (days, title, body) in enumerate(sprints):
        cx = int(T.MARGIN_X) + index * (card_w + int(gap))
        accent = C._mix(T.INDIGO, T.CYAN, index / max(len(sprints) - 1, 1))
        C.diamond(slide, cx + int(card_w / 2), y + Inches(0.11), Inches(0.20), fill=accent)

        cy = y + Inches(0.34)
        C.rect(slide, cx, cy, card_w, card_h, fill=T.SURFACE, line=T.LINE, radius=0.06)
        C.rect(slide, cx, cy, card_w, Inches(0.035), fill=accent)
        pad = Inches(0.18)
        inner_w = card_w - 2 * int(pad)
        C.add_text(slide, cx + int(pad), cy + Inches(0.18), inner_w, Inches(0.22),
                   days.upper(), size=Pt(8.5), color=accent, bold=True, letter_spacing=1.1)
        C.add_text(slide, cx + int(pad), cy + Inches(0.42), inner_w, Inches(0.26),
                   title, size=Pt(13), color=T.TEXT_HI, bold=True)
        size = C.fit_size(body, inner_w, Inches(1.14), 9.5, line_spacing=1.20)
        C.add_text(slide, cx + int(pad), cy + Inches(0.70), inner_w, Inches(1.18), body,
                   size=Pt(size), color=T.TEXT_MID, line_spacing=1.20)

    y += Inches(0.34) + int(card_h) + Inches(0.20)

    phases = data["phases"]
    geo = T.columns(len(phases), gutter=Inches(0.20))
    for index, ((tag, window, scope), (cx, cw)) in enumerate(zip(phases, geo)):
        accent = T.ACCENTS[index % len(T.ACCENTS)]
        height = Inches(0.82)
        C.rect(slide, cx, y, cw, height, fill=T.SURFACE_2, line=T.LINE, radius=0.06)
        C.add_text(slide, int(cx) + Inches(0.18), y + Inches(0.09), Inches(0.60),
                   Inches(0.24), tag, size=Pt(13), color=accent, bold=True)
        C.add_text(slide, int(cx) + Inches(0.80), y + Inches(0.11), Inches(1.0),
                   Inches(0.22), window, size=Pt(9), color=T.TEXT_LOW, bold=True)
        C.add_text(slide, int(cx) + Inches(0.18), y + Inches(0.35),
                   int(cw) - Inches(0.36), Inches(0.40), scope,
                   size=Pt(9.5), color=T.TEXT_MID, line_spacing=1.18)


# ==========================================================================
# Team
# ==========================================================================

_STATUS_COLORS = {"In seat": T.GREEN, "Hiring": T.AMBER, "Contract": T.TEXT_LOW}


def render_team(slide, data: dict, y: int) -> None:
    roles = data["roles"]
    cols = 3
    gap = Inches(0.24)
    card_w = int((int(T.CONTENT_W) - (cols - 1) * int(gap)) / cols)
    card_h = Inches(1.62)

    for index, (role, name, scope, status) in enumerate(roles):
        row, col = divmod(index, cols)
        cx = int(T.MARGIN_X) + col * (card_w + int(gap))
        cy = y + row * (int(card_h) + int(Inches(0.18)))
        accent = _STATUS_COLORS.get(status, T.TEXT_LOW)

        C.rect(slide, cx, cy, card_w, card_h, fill=T.SURFACE, line=T.LINE, radius=0.05)
        C.rect(slide, cx, cy, Inches(0.035), card_h, fill=accent)

        pad = Inches(0.24)
        inner_x = cx + int(pad)
        inner_w = card_w - 2 * int(pad)

        C.chip(slide, inner_x, cy + Inches(0.18), Inches(0.26), status.upper(),
               fill=C._mix(T.INK, accent, 0.20), line=C._mix(T.LINE, accent, 0.45),
               color=accent, size=Pt(8), bold=True)
        C.add_text(slide, inner_x, cy + Inches(0.54), inner_w, Inches(0.26), role,
                   size=Pt(13), color=T.TEXT_HI, bold=True)
        C.add_text(slide, inner_x, cy + Inches(0.82), inner_w, Inches(0.24), name,
                   size=Pt(10.5), color=accent, bold=True)
        C.add_text(slide, inner_x, cy + Inches(1.08), inner_w, Inches(0.42), scope,
                   size=Pt(9.5), color=T.TEXT_MID, line_spacing=1.24)

    rows = (len(roles) + cols - 1) // cols
    y += rows * (int(card_h) + int(Inches(0.18))) + Inches(0.10)
    C.note_bar(slide, T.MARGIN_X, y, T.CONTENT_W, data["note"],
               label="Before you send this", accent=T.AMBER)


# ==========================================================================
# The ask
# ==========================================================================

def render_ask(slide, data: dict, y: int) -> None:
    stats = data["headline_stats"]
    geo = T.columns(len(stats), gutter=Inches(0.24))
    stat_h = Inches(1.14)
    accents = (T.CYAN, T.INDIGO, T.GREEN, T.AMBER)
    for index, ((value, label), (cx, cw)) in enumerate(zip(stats, geo)):
        accent = accents[index % len(accents)]
        C.rect(slide, cx, y, cw, stat_h,
               fill=C._mix(T.SURFACE, accent, 0.12), line=C._mix(T.LINE, accent, 0.38),
               radius=0.05)
        C.rect(slide, cx, y, cw, Inches(0.04), fill=accent)
        size = C.fit_size(value, int(cw) - Inches(0.44), Inches(0.52), 30.0, bold=True,
                          line_spacing=1.0)
        C.add_text(slide, int(cx) + Inches(0.22), y + Inches(0.26),
                   int(cw) - Inches(0.44), Inches(0.48), value,
                   size=Pt(size), color=accent, bold=True, letter_spacing=-0.8)
        C.add_text(slide, int(cx) + Inches(0.22), y + Inches(0.78),
                   int(cw) - Inches(0.44), Inches(0.24), label,
                   size=Pt(10.5), color=T.TEXT_MID)

    y += int(stat_h) + Inches(0.30)

    left_w = Inches(5.55)
    right_x = int(T.MARGIN_X) + int(left_w) + Inches(0.48)
    right_w = int(T.MARGIN_X) + int(T.CONTENT_W) - right_x

    C.add_text(slide, T.MARGIN_X, y, left_w, Inches(0.24),
               data["use_of_funds_label"].upper(), size=T.SIZE_MICRO, color=T.CYAN,
               bold=True, letter_spacing=1.2)
    fy = y + Inches(0.34)
    bar_max = int(left_w) - Inches(2.90)
    largest = max(share for _, share in data["use_of_funds"])
    for index, (label, share) in enumerate(data["use_of_funds"]):
        accent = C._mix(T.INDIGO, T.CYAN, index / max(len(data["use_of_funds"]) - 1, 1))
        C.add_text(slide, T.MARGIN_X, fy, Inches(2.42), Inches(0.24), label,
                   size=Pt(10), color=T.TEXT_MID, anchor=MSO_ANCHOR.MIDDLE)
        track_x = int(T.MARGIN_X) + Inches(2.50)
        C.rect(slide, track_x, fy + Inches(0.055), bar_max, Inches(0.14),
               fill=T.SURFACE_2, radius=0.5)
        C.rect(slide, track_x, fy + Inches(0.055), int(bar_max * share / largest),
               Inches(0.14), fill=accent, radius=0.5)
        C.add_text(slide, int(T.MARGIN_X) + int(left_w) - Inches(0.42), fy,
                   Inches(0.42), Inches(0.24), f"{share}%",
                   size=Pt(10), color=T.TEXT_HI, bold=True, align=PP_ALIGN.RIGHT,
                   anchor=MSO_ANCHOR.MIDDLE)
        fy += Inches(0.32)

    C.add_text(slide, right_x, y, right_w, Inches(0.24),
               data["milestones_label"].upper(), size=T.SIZE_MICRO, color=T.GREEN,
               bold=True, letter_spacing=1.2)
    C.add_bullets(slide, right_x, y + Inches(0.34), right_w, Inches(1.90),
                  data["milestones"], size=Pt(10.5), accent=T.GREEN, space_after=Pt(5))

    y = max(fy, y + Inches(2.30)) + Inches(0.10)
    C.note_bar(slide, T.MARGIN_X, y, T.CONTENT_W, data["assumption"],
               label="Assumption", accent=T.TEXT_LOW, size=T.SIZE_MICRO)


# ==========================================================================
# Dispatch
# ==========================================================================

_DIAGRAMS = {
    "diagram_two_brain": D.two_brain,
    "diagram_layer": D.layer_stack,
    "workflow_story": D.workflow_story,
    "diagram_state_machine": D.state_machine,
    "diagram_agents": D.agent_grid,
    "diagram_pipeline": D.knowledge_pipeline,
    "diagram_security": D.security_stack,
}

_RENDERERS = {
    "cards": render_cards,
    "evidence": render_evidence,
    "autonomy": render_autonomy,
    "wedge": render_wedge,
    "matrix": render_matrix,
    "moat": render_moat,
    "pricing": render_pricing,
    "market": render_market,
    "validation": render_validation,
    "metrics": render_metrics,
    "roadmap": render_roadmap,
    "team": render_team,
    "ask": render_ask,
}


def build() -> Path:
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H

    core = prs.core_properties
    core.title = K.DECK_TITLE
    core.subject = K.DECK_SUBJECT
    core.author = K.DECK_AUTHOR
    core.comments = (
        "Generated by deck/build.py. Edit copy in deck/content.py and rebuild."
    )

    total = len(K.SLIDES)
    for index, data in enumerate(K.SLIDES, start=1):
        kind = data["kind"]
        slide = C.new_slide(prs, dark=True)

        if kind == "cover":
            render_cover(slide, data)
            continue
        if kind == "closing":
            render_closing(slide, data)
            continue
        if kind == "divider":
            render_divider(slide, data)
            C.footer(slide, index=index, total=total)
            continue

        y = C.header(
            slide,
            eyebrow_text=data.get("eyebrow", ""),
            title=data["title"],
            lead=data.get("lead"),
        )

        if kind in _DIAGRAMS:
            _DIAGRAMS[kind](slide, data, y)
        elif kind in _RENDERERS:
            _RENDERERS[kind](slide, data, y)
        else:
            raise ValueError(f"no renderer for slide kind {kind!r}")

        C.footer(slide, index=index, total=total, source=data.get("source", ""))

    T.OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


def main() -> int:
    path = build()
    size_kb = path.stat().st_size / 1024
    print(f"font       {T.FONT}")
    print(f"slides     {len(K.SLIDES)}")
    print(f"written    {path}  ({size_kb:.0f} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
