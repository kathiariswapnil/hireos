"""Reusable slide primitives for the HireOS deck.

Everything here draws native PowerPoint shapes and real text boxes so the
generated deck stays fully editable. Nothing is rasterized except the brand
artwork.

Text measurement is done against the actual installed font via Pillow rather
than guessed, because the layouts pack a lot of copy and silent overflow is the
main failure mode of generated decks.
"""

from __future__ import annotations

import functools
from pathlib import Path

from PIL import ImageFont
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Length, Pt
from lxml import etree

from brand import tokens as T

EMU_PER_IN = 914400
EMU_PER_PT = 12700


# ==========================================================================
# Text measurement
# ==========================================================================

_FONT_FILES = {
    "Inter": [
        ("/Library/Fonts/Inter.ttc", 0),
        (str(Path.home() / "Library/Fonts/Inter.ttc"), 0),
        (str(Path.home() / "Library/Fonts/Inter-Regular.ttf"), 0),
    ],
    "Helvetica Neue": [("/System/Library/Fonts/HelveticaNeue.ttc", 0)],
    "Avenir Next": [("/System/Library/Fonts/Avenir Next.ttc", 0)],
    "Helvetica": [("/System/Library/Fonts/Helvetica.ttc", 0)],
}
_BOLD_INDEX = {"Helvetica Neue": 1, "Helvetica": 1, "Avenir Next": 1}

# Fallback average glyph width as a fraction of point size, used only when no
# font file can be loaded for measurement.
_FALLBACK_RATIO = 0.505


@functools.lru_cache(maxsize=64)
def _pil_font(size_pt: float, bold: bool):
    for path, index in _FONT_FILES.get(T.FONT, []):
        if not Path(path).exists():
            continue
        try:
            idx = _BOLD_INDEX.get(T.FONT, 0) if bold else index
            return ImageFont.truetype(path, max(int(round(size_pt * 4)), 8), index=idx)
        except Exception:
            continue
    return None


def _as_points(size) -> float:
    """Accept either a float point size or a pptx Length and return points.

    Passing Pt(13) where points were expected silently produced shapes tens of
    thousands of inches tall, so the conversion is enforced at the boundary.
    """
    return float(size.pt) if isinstance(size, Length) else float(size)


def text_width_pt(text: str, size_pt: float, bold: bool = False) -> float:
    """Width of a single line of text, in points."""
    size_pt = _as_points(size_pt)
    font = _pil_font(size_pt, bold)
    if font is None:
        return len(text) * size_pt * _FALLBACK_RATIO
    # Measured at 4x the point size for precision, so scale back down.
    return font.getlength(text) / 4.0


def wrap_lines(text: str, width_emu: int, size_pt: float, bold: bool = False) -> list[str]:
    """Greedy word wrap using measured glyph widths."""
    size_pt = _as_points(size_pt)
    limit_pt = width_emu / EMU_PER_PT
    lines: list[str] = []
    for hard_line in text.split("\n"):
        words = hard_line.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            candidate = f"{current} {word}"
            if text_width_pt(candidate, size_pt, bold) <= limit_pt:
                current = candidate
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def text_height(
    text: str,
    width_emu: int,
    size_pt: float,
    *,
    bold: bool = False,
    line_spacing: float = 1.18,
) -> int:
    """Rendered height of wrapped text, in EMU."""
    size_pt = _as_points(size_pt)
    count = max(len(wrap_lines(text, width_emu, size_pt, bold)), 1)
    return int(count * size_pt * line_spacing * EMU_PER_PT)


def fit_size(
    text: str,
    width_emu: int,
    height_emu: int,
    start_pt: float,
    *,
    bold: bool = False,
    line_spacing: float = 1.18,
    min_pt: float = 7.0,
) -> float:
    """Largest point size at or below start_pt where the text still fits."""
    size = _as_points(start_pt)
    while size > min_pt:
        if text_height(text, width_emu, size, bold=bold, line_spacing=line_spacing) <= height_emu:
            return size
        size -= 0.5
    return min_pt


# ==========================================================================
# Low-level XML helpers
# ==========================================================================

_PPR_ORDER = [
    qn(t)
    for t in (
        "a:lnSpc", "a:spcBef", "a:spcAft",
        "a:buClrTx", "a:buClr",
        "a:buSzTx", "a:buSzPct", "a:buSzPts",
        "a:buFontTx", "a:buFont",
        "a:buNone", "a:buAutoNum", "a:buChar", "a:buBlip",
        "a:tabLst", "a:defRPr", "a:extLst",
    )
]


def _insert_ordered(pPr, element) -> None:
    """Insert a child into a:pPr at its schema-valid position.

    Order matters here: PowerPoint rejects the file outright and shows a repair
    dialog if pPr children appear out of sequence.
    """
    position = _PPR_ORDER.index(element.tag)
    for child in pPr:
        if child.tag in _PPR_ORDER and _PPR_ORDER.index(child.tag) > position:
            child.addprevious(element)
            return
    pPr.append(element)


def _el(tag: str, **attrs):
    element = etree.SubElement(etree.Element(qn("a:dummy")), qn(tag))
    for key, value in attrs.items():
        element.set(key, value)
    return element


def set_bullet(paragraph, char: str, color: RGBColor, size_pct: int = 100) -> None:
    """Apply a native PowerPoint character bullet with a hanging indent."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(Inches(0.19)))
    pPr.set("indent", str(-Inches(0.19)))

    bu_clr = _el("a:buClr")
    srgb = etree.SubElement(bu_clr, qn("a:srgbClr"))
    srgb.set("val", str(color))
    _insert_ordered(pPr, bu_clr)

    _insert_ordered(pPr, _el("a:buSzPct", val=str(int(size_pct * 1000))))
    _insert_ordered(pPr, _el("a:buFont", typeface="Arial"))
    _insert_ordered(pPr, _el("a:buChar", char=char))


def no_shadow(shape) -> None:
    """Strip the inherited default drop shadow from an autoshape."""
    try:
        shape.shadow.inherit = False
    except (AttributeError, NotImplementedError):
        pass


# ==========================================================================
# Slide + background
# ==========================================================================

def new_slide(prs, *, dark: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = T.INK if dark else T.PAPER
    return slide


def add_picture_cover(slide, path: Path, left, top, width, height):
    """Place a picture, center-cropping it to exactly fill the target box."""
    from PIL import Image

    with Image.open(path) as image:
        src_w, src_h = image.size
    box_ratio = width / height
    src_ratio = src_w / src_h

    if src_ratio > box_ratio:
        draw_h = height
        draw_w = int(height * src_ratio)
    else:
        draw_w = width
        draw_h = int(width / src_ratio)

    picture = slide.shapes.add_picture(
        str(path), Emu(int(left)), Emu(int(top)), Emu(draw_w), Emu(draw_h)
    )
    # Crop the overflow symmetrically so the art stays centered in the box.
    if draw_w > width:
        overflow = (draw_w - width) / draw_w / 2
        picture.crop_left = overflow
        picture.crop_right = overflow
    if draw_h > height:
        overflow = (draw_h - height) / draw_h / 2
        picture.crop_top = overflow
        picture.crop_bottom = overflow
    picture.left, picture.top = Emu(int(left)), Emu(int(top))
    picture.width, picture.height = Emu(int(width)), Emu(int(height))
    return picture


# ==========================================================================
# Shapes
# ==========================================================================

def set_picture_transparency(picture, amount: float) -> None:
    """Fade a picture by `amount` (0 = untouched, 1 = invisible).

    Used to sink decorative art behind text without a separate exported asset.
    """
    blip = picture._element.blipFill.find(qn("a:blip"))
    if blip is None:
        return
    existing = blip.find(qn("a:alphaModFix"))
    if existing is not None:
        blip.remove(existing)
    alpha = etree.SubElement(blip, qn("a:alphaModFix"))
    alpha.set("amt", str(int((1.0 - amount) * 100000)))


def rect(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill=None,
    line=None,
    line_w=Pt(1),
    radius: float | None = None,
    shape_type=None,
):
    if shape_type is None:
        shape_type = (
            MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
        )
    shape = slide.shapes.add_shape(
        shape_type, Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    if radius is not None and shape.adjustments:
        shape.adjustments[0] = radius
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_w
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def gradient_bar(slide, left, top, width, height, start, end, angle: float = 0.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    shape.line.fill.background()
    shape.shadow.inherit = False
    fill = shape.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = start
    stops[0].position = 0.0
    stops[1].color.rgb = end
    stops[1].position = 1.0
    fill.gradient_angle = angle
    return shape


def _set_gradient_alpha(shape, alphas: tuple[float, ...]) -> None:
    """Apply per-stop alpha to a gradient fill.

    python-pptx exposes gradient stop colors but not their transparency, so the
    alpha element is written directly onto each stop's color.
    """
    stops = list(shape._element.iter(qn("a:gs")))
    for stop, alpha in zip(stops, alphas):
        srgb = stop.find(qn("a:srgbClr"))
        if srgb is None:
            continue
        existing = srgb.find(qn("a:alpha"))
        if existing is not None:
            srgb.remove(existing)
        element = etree.SubElement(srgb, qn("a:alpha"))
        element.set("val", str(int(max(0.0, min(1.0, alpha)) * 100000)))


def scrim(
    slide,
    left,
    top,
    width,
    height,
    *,
    color=None,
    angle: float = 0.0,
    start_alpha: float = 1.0,
    end_alpha: float = 0.0,
):
    """A gradient fade used to keep text legible over full-bleed artwork."""
    color = T.INK if color is None else color
    shape = gradient_bar(slide, left, top, width, height, color, color, angle)
    _set_gradient_alpha(shape, (start_alpha, end_alpha))
    return shape


def diamond(slide, cx, cy, size, *, fill=None, line=None, line_w=Pt(1.25)):
    return rect(
        slide,
        cx - size / 2,
        cy - size / 2,
        size,
        size,
        fill=fill,
        line=line,
        line_w=line_w,
        shape_type=MSO_SHAPE.DIAMOND,
    )


def connector(slide, x1, y1, x2, y2, color, width=Pt(1)):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2))
    )
    line.line.color.rgb = color
    line.line.width = width
    return line


def arrow(slide, x1, y1, x2, y2, color, width=Pt(1)):
    line = connector(slide, x1, y1, x2, y2, color, width)
    line_el = line.line._get_or_add_ln()
    tail = line_el.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(line_el, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    return line


# ==========================================================================
# Text
# ==========================================================================

def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size=T.SIZE_BODY,
    color=T.TEXT_MID,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing: float = 1.18,
    space_after=Pt(0),
    font=None,
    letter_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(
        Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor

    for index, chunk in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = space_after
        run = paragraph.add_run()
        run.text = chunk
        run.font.size = size
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font or T.FONT
        run.font.color.rgb = color
        if letter_spacing is not None:
            run.font._rPr.set("spc", str(int(letter_spacing * 100)))
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    *,
    size=T.SIZE_BODY,
    color=T.TEXT_MID,
    accent=T.INDIGO,
    line_spacing: float = 1.2,
    space_after=Pt(6),
    bullet_char: str = "\u2022",
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(
        Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.line_spacing = line_spacing
        paragraph.space_after = space_after
        run = paragraph.add_run()
        run.text = item
        run.font.size = size
        run.font.name = T.FONT
        run.font.color.rgb = color
        set_bullet(paragraph, bullet_char, accent, size_pct=90)
    return box


def bullets_height(
    items: list[str],
    width_emu: int,
    size_pt: float,
    *,
    line_spacing: float = 1.2,
    space_after_pt: float = 6.0,
) -> int:
    """Height of an add_bullets block, accounting for the hanging indent."""
    inner = width_emu - Inches(0.19)
    total = 0
    for item in items:
        lines = max(len(wrap_lines(item, inner, size_pt)), 1)
        total += int(lines * size_pt * line_spacing * EMU_PER_PT)
        total += int(space_after_pt * EMU_PER_PT)
    return total


# ==========================================================================
# Slide chrome
# ==========================================================================

def eyebrow(slide, text: str, *, color=T.CYAN, top=None, left=None, dark=True):
    top = T.MARGIN_TOP if top is None else top
    left = T.MARGIN_X if left is None else left
    tick = rect(slide, left, int(top) + Inches(0.035), Inches(0.115), Inches(0.115),
                fill=color, shape_type=MSO_SHAPE.DIAMOND)
    add_text(
        slide,
        int(left) + Inches(0.24),
        top,
        Inches(8.0),
        Inches(0.22),
        text.upper(),
        size=T.SIZE_EYEBROW,
        color=color,
        bold=True,
        letter_spacing=1.4,
    )
    return tick


def header(
    slide,
    *,
    eyebrow_text: str,
    title: str,
    lead: str | None = None,
    dark: bool = True,
    lead_width=None,
    title_size=None,
) -> int:
    """Draw eyebrow + title + optional lead. Returns the content-start Y in EMU."""
    title_color = T.TEXT_HI if dark else T.INK_TEXT_HI
    lead_color = T.TEXT_MID if dark else T.INK_TEXT_MID
    accent = T.CYAN if dark else T.INDIGO_DEEP

    if eyebrow_text:
        eyebrow(slide, eyebrow_text, color=accent, dark=dark)

    y = int(T.MARGIN_TOP) + Inches(0.40)
    title_width = int(T.CONTENT_W * 0.86)
    size_pt = title_size if title_size else T.SIZE_TITLE.pt
    size_pt = fit_size(title, title_width, Inches(1.05), size_pt, bold=True, line_spacing=1.1)
    t_height = text_height(title, title_width, size_pt, bold=True, line_spacing=1.1)
    add_text(
        slide, T.MARGIN_X, y, title_width, t_height + Inches(0.06), title,
        size=Pt(size_pt), color=title_color, bold=True, line_spacing=1.1,
    )
    y += t_height + Inches(0.16)

    if lead:
        width = lead_width if lead_width else int(T.CONTENT_W * 0.80)
        l_height = text_height(lead, width, T.SIZE_LEAD.pt, line_spacing=1.3)
        add_text(
            slide, T.MARGIN_X, y, width, l_height + Inches(0.06), lead,
            size=T.SIZE_LEAD, color=lead_color, line_spacing=1.3,
        )
        y += l_height + Inches(0.30)
    else:
        y += Inches(0.10)

    return int(y)


def footer(slide, *, index: int, total: int, source: str = "", dark: bool = True):
    y = int(T.SLIDE_H) - int(T.MARGIN_BOTTOM)
    muted = T.TEXT_LOW if dark else T.INK_TEXT_LOW

    rect(slide, T.MARGIN_X, y - Inches(0.10), T.CONTENT_W, Emu(9525),
         fill=T.LINE if dark else T.PAPER_2)

    add_text(
        slide, T.MARGIN_X, y, Inches(3.0), Inches(0.2),
        T.BRAND_NAME, size=T.SIZE_MICRO, color=muted, bold=True, letter_spacing=1.0,
    )
    if source:
        add_text(
            slide, int(T.MARGIN_X) + Inches(1.2), y, Inches(7.6), Inches(0.2),
            source, size=T.SIZE_MICRO, color=muted,
        )
    add_text(
        slide, int(T.SLIDE_W) - int(T.MARGIN_X) - Inches(1.2), y, Inches(1.2), Inches(0.2),
        f"{index:02d} / {total:02d}", size=T.SIZE_MICRO, color=muted, align=PP_ALIGN.RIGHT,
    )


def logo_lockup(slide, left, top, mark_size, *, dark: bool = True, show_name: bool = True):
    path = T.LOGO_MARK_DARK if dark else T.LOGO_MARK_LIGHT
    if path.exists():
        slide.shapes.add_picture(
            str(path), Emu(int(left)), Emu(int(top)), Emu(int(mark_size)), Emu(int(mark_size))
        )
    if not show_name:
        return
    add_text(
        slide,
        int(left) + int(mark_size) + Inches(0.16),
        int(top) + int(mark_size * 0.16),
        Inches(3.2),
        Emu(int(mark_size * 0.72)),
        T.BRAND_NAME,
        size=Pt(mark_size / EMU_PER_IN * 46),
        color=T.TEXT_HI if dark else T.INK_TEXT_HI,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
        letter_spacing=-0.4,
    )


# ==========================================================================
# Composite blocks
# ==========================================================================

def card(
    slide,
    left,
    top,
    width,
    height,
    *,
    heading: str,
    body: str = "",
    accent=T.INDIGO,
    dark: bool = True,
    heading_size=Pt(15),
    body_size=T.SIZE_SMALL,
    number: str = "",
):
    surface = T.SURFACE if dark else RGBColor(0xFF, 0xFF, 0xFF)
    border = T.LINE if dark else T.PAPER_2
    head_color = T.TEXT_HI if dark else T.INK_TEXT_HI
    body_color = T.TEXT_MID if dark else T.INK_TEXT_MID

    rect(slide, left, top, width, height, fill=surface, line=border, radius=T.RADIUS_ADJ)
    rect(slide, left, top, Inches(0.035), height, fill=accent)

    pad = Inches(0.26)
    inner_left = int(left) + int(pad)
    inner_width = int(width) - 2 * int(pad)
    y = int(top) + int(Inches(0.22))

    if number:
        add_text(
            slide, inner_left, y, inner_width, Inches(0.24), number,
            size=T.SIZE_MICRO, color=accent, bold=True, letter_spacing=1.2,
        )
        y += Inches(0.28)

    h_size = fit_size(heading, inner_width, Inches(0.62), heading_size.pt, bold=True, line_spacing=1.12)
    h_height = text_height(heading, inner_width, h_size, bold=True, line_spacing=1.12)
    add_text(
        slide, inner_left, y, inner_width, h_height + Inches(0.04), heading,
        size=Pt(h_size), color=head_color, bold=True, line_spacing=1.12,
    )
    y += h_height + Inches(0.13)

    if body:
        remaining = int(top) + int(height) - y - int(Inches(0.20))
        b_size = fit_size(body, inner_width, remaining, body_size.pt, line_spacing=1.3)
        add_text(
            slide, inner_left, y, inner_width, remaining, body,
            size=Pt(b_size), color=body_color, line_spacing=1.3,
        )


def chip(
    slide,
    left,
    top,
    height,
    text: str,
    *,
    fill=None,
    line=T.LINE,
    color=T.TEXT_MID,
    size=T.SIZE_SMALL,
    pad=Inches(0.16),
    bold=False,
    min_width=None,
):
    width = int(text_width_pt(text, size.pt, bold) * EMU_PER_PT) + 2 * int(pad)
    if min_width:
        width = max(width, int(min_width))
    rect(slide, left, top, width, height, fill=fill, line=line, radius=0.28)
    add_text(
        slide, left, top, width, height, text,
        size=size, color=color, bold=bold, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    return width


def stat(
    slide,
    left,
    top,
    width,
    *,
    value: str,
    label: str,
    accent=T.CYAN,
    dark: bool = True,
    value_size=T.SIZE_STAT,
):
    size_pt = fit_size(value, width, Inches(0.72), value_size.pt, bold=True, line_spacing=1.0)
    add_text(
        slide, left, top, width, Inches(0.66), value,
        size=Pt(size_pt), color=accent, bold=True, line_spacing=1.0, letter_spacing=-0.6,
    )
    add_text(
        slide, left, int(top) + Inches(0.70), width, Inches(0.48), label,
        size=T.SIZE_SMALL, color=T.TEXT_MID if dark else T.INK_TEXT_MID, line_spacing=1.22,
    )


def styled_table(
    slide,
    left,
    top,
    width,
    height,
    columns: list[str],
    rows: list[tuple],
    *,
    col_widths: list[float] | None = None,
    highlight_column: int | None = None,
    dark: bool = True,
    header_size=T.SIZE_SMALL,
    body_size=T.SIZE_SMALL,
    row_height=None,
):
    shape = slide.shapes.add_table(
        len(rows) + 1, len(columns),
        Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)),
    )
    table = shape.table
    table.first_row = True
    # Banding is applied explicitly below; the theme's own banding fights it.
    table.horz_banding = False

    if col_widths:
        total = sum(col_widths)
        for index, share in enumerate(col_widths):
            table.columns[index].width = Emu(int(width * share / total))

    header_h = int(row_height or Inches(0.42))
    table.rows[0].height = Emu(header_h)
    body_h = int((int(height) - header_h) / max(len(rows), 1))
    for index in range(1, len(rows) + 1):
        table.rows[index].height = Emu(body_h)

    surface = T.SURFACE if dark else RGBColor(0xFF, 0xFF, 0xFF)
    alt = T.SURFACE_2 if dark else T.PAPER
    border = T.LINE if dark else T.PAPER_2
    head_bg = T.SURFACE_2 if dark else T.PAPER_2
    head_fg = T.TEXT_HI if dark else T.INK_TEXT_HI
    body_fg = T.TEXT_MID if dark else T.INK_TEXT_MID

    def style_cell(cell, text, *, bg, fg, bold, size, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.margin_left = cell.margin_right = Inches(0.14)
        cell.margin_top = cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame = cell.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        paragraph.line_spacing = 1.15
        run = paragraph.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.name = T.FONT
        run.font.color.rgb = fg
        _cell_borders(cell, border)

    for index, label in enumerate(columns):
        is_highlight = index == highlight_column
        style_cell(
            table.cell(0, index), label,
            bg=T.INDIGO_DEEP if is_highlight else head_bg,
            fg=T.TEXT_HI if is_highlight else head_fg,
            bold=True, size=header_size,
        )

    for r_index, row in enumerate(rows):
        base_bg = surface if r_index % 2 == 0 else alt
        for c_index, value in enumerate(row):
            is_highlight = c_index == highlight_column
            style_cell(
                table.cell(r_index + 1, c_index), str(value),
                bg=_mix(base_bg, T.INDIGO_DEEP, 0.22) if is_highlight else base_bg,
                fg=T.TEXT_HI if (is_highlight or c_index == 0) else body_fg,
                bold=is_highlight or c_index == 0,
                size=body_size,
            )
    return table


def _cell_borders(cell, color) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tc_pr.find(qn(tag))
        if existing is not None:
            tc_pr.remove(existing)
    # tcPr children must appear in schema order: lnL, lnR, lnT, lnB, then fill.
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):
        line = etree.Element(qn(tag))
        line.set("w", str(int(Pt(0.75))))
        line.set("cap", "flat")
        line.set("cmpd", "sng")
        line.set("algn", "ctr")
        fill = etree.SubElement(line, qn("a:solidFill"))
        srgb = etree.SubElement(fill, qn("a:srgbClr"))
        srgb.set("val", str(color))
        tc_pr.insert(0, line)


def _mix(a: RGBColor, b: RGBColor, ratio: float) -> RGBColor:
    return RGBColor(
        int(a[0] + (b[0] - a[0]) * ratio),
        int(a[1] + (b[1] - a[1]) * ratio),
        int(a[2] + (b[2] - a[2]) * ratio),
    )


def note_bar(
    slide,
    left,
    top,
    width,
    text: str,
    *,
    accent=T.AMBER,
    label: str = "",
    dark: bool = True,
    size=T.SIZE_SMALL,
):
    """A labelled callout used for assumptions, rules and disclaimers."""
    pad = Inches(0.20)
    inner_width = int(width) - 2 * int(pad) - int(Inches(0.06))
    body = f"{label}  {text}" if label else text
    height = text_height(body, inner_width, size.pt, line_spacing=1.3) + 2 * int(pad)

    rect(slide, left, top, width, height,
         fill=_mix(T.SURFACE, accent, 0.10) if dark else _mix(T.PAPER, accent, 0.12),
         line=_mix(T.LINE, accent, 0.35) if dark else _mix(T.PAPER_2, accent, 0.30),
         radius=0.10)
    rect(slide, left, top, Inches(0.045), height, fill=accent)

    box = slide.shapes.add_textbox(
        Emu(int(left) + int(pad) + Inches(0.06)), Emu(int(top) + int(pad)),
        Emu(inner_width), Emu(height - 2 * int(pad)),
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.line_spacing = 1.3
    if label:
        run = paragraph.add_run()
        run.text = f"{label}  "
        run.font.size = size
        run.font.bold = True
        run.font.name = T.FONT
        run.font.color.rgb = accent
    run = paragraph.add_run()
    run.text = text
    run.font.size = size
    run.font.name = T.FONT
    run.font.color.rgb = T.TEXT_MID if dark else T.INK_TEXT_MID
    return height


def kicker(slide, left, top, width, text: str, *, accent=T.CYAN, dark: bool = True):
    inner = int(width) - int(Inches(0.30))
    height = text_height(text, inner, 13.0, bold=True, line_spacing=1.28)
    rect(slide, left, top, Inches(0.045), height + Inches(0.06), fill=accent)
    add_text(
        slide, int(left) + Inches(0.30), top, inner, height + Inches(0.06), text,
        size=Pt(13), color=T.TEXT_HI if dark else T.INK_TEXT_HI, bold=True, line_spacing=1.28,
    )
    return height + Inches(0.06)
